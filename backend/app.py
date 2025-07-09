# 🚀 GENERADOR DE PRESENTACIONES EDUCATIVAS - BACKEND
# ===================================================
# 
# 📋 DESCRIPCIÓN GENERAL:
# Este backend proporciona una API REST para el generador de presentaciones.
# Maneja la generación de contenido con IA, obtención de imágenes y procesamiento de archivos PPTX.
#
# 🎯 FUNCIONALIDADES PRINCIPALES:
# 1. Generar contenido educativo usando OpenAI API
# 2. Obtener imágenes relacionadas usando Unsplash API  
# 3. Leer y procesar archivos PowerPoint (.pptx) subidos por el usuario
# 4. Convertir diapositivas a formato HTML para vista previa
#
# 🔧 DEPENDENCIAS REQUERIDAS:
# pip install flask flask-cors python-pptx requests python-dotenv
#
# 🔑 CONFIGURACIÓN:
# Crear archivo .env con las API keys:
# OPENAI_API_KEY=tu_api_key_de_openai
# UNSPLASH_API_KEY=tu_api_key_de_unsplash

# ===================================================
# 📦 IMPORTACIÓN DE LIBRERÍAS
# ===================================================

from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import requests
from dotenv import load_dotenv

# 🔑 Cargar variables de entorno desde archivo .env
# Esto permite usar API keys de forma segura sin exponerlas en el código
load_dotenv()

# 📄 Importar librería para leer archivos PowerPoint
try:
    from pptx import Presentation
except ModuleNotFoundError:
    raise ImportError("🚨 Necesitas instalar 'python-pptx'. Usa: pip install python-pptx")

# ===================================================
# ⚙️ CONFIGURACIÓN INICIAL
# ===================================================

# 🌐 Crear aplicación Flask
app = Flask(__name__)

# 🔓 Habilitar CORS para permitir peticiones desde el frontend
# Sin esto, el navegador bloquearía las peticiones por seguridad
CORS(app)

# 🔑 Cargar API keys desde variables de entorno
# Si no están configuradas, serán None y se mostrará un error
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
UNSPLASH_API_KEY = os.getenv('UNSPLASH_API_KEY')

# ===================================================
# 📄 ENDPOINT: PROCESAR ARCHIVOS POWERPOINT (.PPTX)
# ===================================================
# 
# 🎯 FUNCIÓN: Leer archivos PowerPoint subidos por el usuario y extraer su contenido
# 📥 ENTRADA: Archivo .pptx enviado desde el frontend
# 📤 SALIDA: Array de diapositivas en formato HTML
# 🔗 URL: POST /subir-pptx

@app.route('/subir-pptx', methods=['POST'])
def subir_pptx():
    """
    📄 Procesa archivos PowerPoint (.pptx) subidos por el usuario
    
    Pasos del proceso:
    1. Recibe el archivo desde el frontend
    2. Valida que sea un archivo .pptx válido
    3. Lee cada diapositiva usando python-pptx
    4. Extrae el texto de cada elemento de la diapositiva
    5. Convierte todo a formato HTML para mostrar en el frontend
    """
    
    # 📁 Obtener el archivo enviado desde el frontend
    file = request.files.get('archivo')

    # ✅ VALIDACIÓN: Verificar que el archivo existe y es un .pptx
    if not file or not file.filename or not file.filename.endswith('.pptx'):
        return jsonify({'error': 'Archivo .pptx no válido o no proporcionado'}), 400

    try:
        # 📖 Abrir el archivo PowerPoint usando python-pptx
        # file.stream proporciona el contenido del archivo como bytes
        prs = Presentation(file.stream)
    except Exception as e:
        # ❌ Si hay error al leer el archivo, devolver error 500
        return jsonify({'error': f'Error al leer el archivo: {str(e)}'}), 500

    # 📝 Array para almacenar las diapositivas convertidas a HTML
    slides_html = []

    # 🔄 ITERAR POR CADA DIAPOSITIVA DEL POWERPOINT
    for idx, slide in enumerate(prs.slides):
        contenido = ''
        
        # 🔍 BUSCAR TEXTO EN CADA ELEMENTO DE LA DIAPOSITIVA
        for shape in slide.shapes:
            try:
                # ✅ Verificar si el elemento tiene texto
                if hasattr(shape, 'text_frame') and shape.text_frame:  # type: ignore
                    # 📝 Extraer el texto y agregarlo al contenido
                    contenido += f"<p>{shape.text_frame.text}</p>"  # type: ignore
            except (AttributeError, TypeError):
                # ⚠️ Si el elemento no tiene texto, continuar con el siguiente
                continue

        # 🧱 CREAR HTML PARA ESTA DIAPOSITIVA
        # Incluye el número de diapositiva y el contenido extraído
        slide_html = f"""
        <div class='slide'>
            <h2>Diapositiva {idx + 1}</h2>
            <div class='slide-content'>
                {contenido if contenido else '<em>(sin texto)</em>'}
            </div>
        </div>"""

        # ➕ Agregar la diapositiva al array de resultados
        slides_html.append(slide_html)

    # 📤 DEVOLVER TODAS LAS DIAPOSITIVAS EN FORMATO JSON
    return jsonify({"slides": slides_html})

# ===================================================
# 🤖 ENDPOINT: GENERAR CONTENIDO CON OPENAI API
# ===================================================
# 
# 🎯 FUNCIÓN: Generar contenido educativo usando inteligencia artificial
# 📥 ENTRADA: Título del tema a explicar
# 📤 SALIDA: Texto explicativo generado por IA
# 🔗 URL: POST /generar-contenido
# 🔑 SEGURIDAD: Usa API key desde variables de entorno

@app.route('/generar-contenido', methods=['POST'])
def generar_contenido():
    """
    🤖 Genera contenido educativo usando OpenAI API
    
    Este endpoint:
    1. Recibe un título de tema desde el frontend
    2. Envía una petición a OpenAI con un prompt educativo
    3. Recibe texto explicativo generado por IA
    4. Devuelve el contenido al frontend
    
    El prompt está diseñado para generar explicaciones claras y educativas
    en 3 frases, perfectas para presentaciones.
    """
    
    # 📥 Obtener datos JSON enviados desde el frontend
    data = request.json or {}
    titulo = data.get('titulo')
    
    # 🔑 VALIDACIÓN: Verificar que la API key esté configurada
    if not OPENAI_API_KEY:
        return jsonify({'error': 'API key de OpenAI no configurada'}), 500
    
    try:
        # 🌐 HACER PETICIÓN A OPENAI API
        # Usamos GPT-3.5-turbo para generar contenido educativo
        response = requests.post('https://api.openai.com/v1/chat/completions', 
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {OPENAI_API_KEY}'  # 🔑 API key segura
            },
            json={
                'model': 'gpt-3.5-turbo',  # 🤖 Modelo de IA a usar
                'messages': [{
                    'role': 'user', 
                    'content': f'Explica este tema en 3 frases claras: {titulo}'
                }]
            }
        )
        
        # ✅ VERIFICAR SI LA PETICIÓN FUE EXITOSA
        if response.status_code == 200:
            # 📝 Extraer el texto generado por la IA
            data = response.json()
            texto_generado = data['choices'][0]['message']['content']
            return jsonify({'texto': texto_generado})
        else:
            # ❌ Si OpenAI devuelve error, informar al frontend
            return jsonify({'error': 'Error al generar contenido'}), 500
            
    except Exception as e:
        # ❌ Si hay error de conexión u otro problema
        return jsonify({'error': str(e)}), 500

# ===================================================
# 🖼️ ENDPOINT: OBTENER IMÁGENES DE UNSPLASH
# ===================================================
# 
# 🎯 FUNCIÓN: Obtener imágenes relacionadas con el tema de la presentación
# 📥 ENTRADA: Palabra clave para buscar imágenes
# 📤 SALIDA: URL de imagen de alta calidad
# 🔗 URL: POST /obtener-imagen
# 🔑 SEGURIDAD: Usa API key desde variables de entorno

@app.route('/obtener-imagen', methods=['POST'])
def obtener_imagen():
    """
    🖼️ Obtiene imágenes relacionadas usando Unsplash API
    
    Este endpoint:
    1. Recibe una palabra clave desde el frontend
    2. Busca imágenes relacionadas en Unsplash
    3. Devuelve la URL de una imagen aleatoria de alta calidad
    4. Las imágenes son gratuitas y de uso libre
    
    Unsplash proporciona imágenes profesionales perfectas para presentaciones.
    """
    
    # 📥 Obtener datos JSON enviados desde el frontend
    data = request.json or {}
    query = data.get('query')
    
    # 🔑 VALIDACIÓN: Verificar que la API key esté configurada
    if not UNSPLASH_API_KEY:
        return jsonify({'error': 'API key de Unsplash no configurada'}), 500
    
    try:
        # 🌐 HACER PETICIÓN A UNSPLASH API
        # Buscamos una imagen aleatoria relacionada con el tema
        response = requests.get(f'https://api.unsplash.com/photos/random?query={query}', 
            headers={'Authorization': f'Client-ID {UNSPLASH_API_KEY}'}  # 🔑 API key segura
        )
        
        # ✅ VERIFICAR SI LA PETICIÓN FUE EXITOSA
        if response.status_code == 200:
            # 🖼️ Extraer la URL de la imagen
            data = response.json()
            url_imagen = data['urls']['regular']  # Imagen de calidad media-alta
            return jsonify({'imagen': url_imagen})
        else:
            # ❌ Si Unsplash devuelve error, informar al frontend
            return jsonify({'error': 'Error al obtener imagen'}), 500
            
    except Exception as e:
        # ❌ Si hay error de conexión u otro problema
        return jsonify({'error': str(e)}), 500

# ===================================================
# 🚀 INICIAR EL SERVIDOR
# ===================================================
# 
# 🎯 FUNCIÓN: Punto de entrada para ejecutar el servidor Flask
# 📍 PUERTO: 5000 (http://localhost:5000)
# 🔧 MODO: Debug activado para desarrollo
# 📝 LOGS: Muestra información detallada de peticiones y errores

if __name__ == '__main__':
    print("🚀 Iniciando servidor del Generador de Presentaciones...")
    print("📍 Servidor disponible en: http://localhost:5000")
    print("🔧 Modo debug: ACTIVADO")
    print("📋 Endpoints disponibles:")
    print("   - POST /subir-pptx (procesar archivos PowerPoint)")
    print("   - POST /generar-contenido (generar texto con OpenAI)")
    print("   - POST /obtener-imagen (obtener imágenes de Unsplash)")
    print("=" * 50)
    
    # 🌐 Iniciar el servidor Flask
    # debug=True: Muestra errores detallados y recarga automáticamente
    # port=5000: Puerto estándar para desarrollo
    app.run(debug=True, port=5000)
